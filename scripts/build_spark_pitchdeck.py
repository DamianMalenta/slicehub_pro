#!/usr/bin/env python3
"""
Build Spark 3.0 Accelerator pitch deck (PPTX + PDF) from repo assets + canonical docs pearls.

Assets used (committed in repo):
  - modules/online/screenshots/wide.svg, narrow.svg  (PWA storefront mock)
  - modules/pos/screenshots/wide.svg, narrow.svg    (PWA POS mock)
  - modules/online/icon.svg                           (optional small logo)

Outputs:
  - _docs/pitchdeck_SPARK_3_assets/png/*.png          (rasterized SVG, 2x scale)
  - _docs/pitchdeck_SPARK_3_output/SliceHub_Spark3_Pitch.pptx
  - _docs/pitchdeck_SPARK_3_output/SliceHub_Spark3_Pitch.pdf

Requires: pip install python-pptx pillow cairosvg fpdf2
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "_docs" / "pitchdeck_SPARK_3_output"
PNG_DIR = ROOT / "_docs" / "pitchdeck_SPARK_3_assets" / "png"
SVG_SOURCES = [
    ("online_wide", ROOT / "modules" / "online" / "screenshots" / "wide.svg"),
    ("online_narrow", ROOT / "modules" / "online" / "screenshots" / "narrow.svg"),
    ("pos_wide", ROOT / "modules" / "pos" / "screenshots" / "wide.svg"),
    ("pos_narrow", ROOT / "modules" / "pos" / "screenshots" / "narrow.svg"),
]


def rasterize_svgs() -> dict[str, Path]:
    import cairosvg

    PNG_DIR.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}
    for key, svg_path in SVG_SOURCES:
        if not svg_path.exists():
            print(f"WARN: missing {svg_path}", file=sys.stderr)
            continue
        out = PNG_DIR / f"{key}.png"
        cairosvg.svg2png(url=str(svg_path), write_to=str(out), scale=2.0)
        paths[key] = out
        print(f"OK raster {key} -> {out}")
    return paths


def build_pptx(png: dict[str, Path]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title: str, subtitle: str) -> None:
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12), Inches(1.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(244, 231, 210)
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(243, 199, 109)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(10, 8, 6)

    def add_dark_blank():
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(16, 14, 12)
        return slide

    def add_bullets(title: str, lines: list[str], foot: str | None = None):
        slide = add_dark_blank()
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.9))
        tfp = tb.text_frame
        tfp.paragraphs[0].text = title
        tfp.paragraphs[0].font.size = Pt(32)
        tfp.paragraphs[0].font.bold = True
        tfp.paragraphs[0].font.color.rgb = RGBColor(243, 199, 109)

        body = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(12.2), Inches(5.5))
        bf = body.text_frame
        bf.word_wrap = True
        for i, line in enumerate(lines):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = line
            para.font.size = Pt(15)
            para.font.color.rgb = RGBColor(230, 225, 215)
            para.space_after = Pt(10)
        if foot:
            fp = bf.add_paragraph()
            fp.text = foot
            fp.font.size = Pt(11)
            fp.font.italic = True
            fp.font.color.rgb = RGBColor(160, 150, 140)

    def add_image_slide(title: str, png_path: Path, caption: str):
        slide = add_dark_blank()
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.75))
        tb.text_frame.paragraphs[0].text = title
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(243, 199, 109)
        slide.shapes.add_picture(str(png_path), Inches(0.55), Inches(1.15), width=Inches(12.2))
        cb = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5))
        cb.text_frame.paragraphs[0].text = caption
        cb.text_frame.paragraphs[0].font.size = Pt(11)
        cb.text_frame.paragraphs[0].font.color.rgb = RGBColor(170, 160, 150)

    # --- slides ---
    add_title_slide(
        "SliceHub Enterprise",
        "System operacyjny gastronomii · Spark 3.0 Accelerator\n"
        "„Okno do restauracji, nie siatka produktów” (_docs/00_PAMIEC_SYSTEMU — Prawo VII)",
    )

    add_bullets(
        "Problem rynku",
        [
            "Lokal utyka między agregatorem (ginie marka) a backoffice (POS, magazyn, ceny) bez spójnego doświadczenia online.",
            "Klient widzi „sklep SKU”; restauracja traci teatr marki i nie spiętego food costu z tym, co wisi na zdjęciu.",
            "„AI w menu” często = kolejny filtr — nie przewaga produktowa (nasze Prawo VII: innowacja albo nic).",
        ],
        "Źródło: _docs/00_PAMIEC_SYSTEMU.md, _docs/canvasy/SliceHub Three Directions.md",
    )

    add_bullets(
        "Rozwiązanie — jedna platforma",
        [
            "Wielonajemczy OS: menu temporalne, macierz cen omnichannel (POS / Takeaway / Delivery), magazyn, logistyka, POS, KDS, kelner, stoliki.",
            "Manager: Studio + Director (scena dania) + Style Conductor + Harmony Score — „kinowe zaplecze”, nie tylko CRUD.",
            "Klient: Scena Drzwi → Counter / Living Table → checkout gościnny → Track z ETA i mapą kierowcy.",
        ],
        "Źródło: _docs/15_KIERUNEK_ONLINE.md (droga B: Counter + Drzwi), _docs/00_PAMIEC_SYSTEMU.md",
    )

    add_bullets(
        "Perły techniczne — integralność i bezpieczeństwo",
        [
            "CartEngine wyłącznie na serwerze — frontend nigdy nie wysyła cen ani totali (Prawo IV).",
            "Multi-tenant: każde SQL z barierą tenant_id; silosy sh_ / sys_ / wh_ łączone tylko po SKU / ascii_key (Konstytucja §9).",
            "Event outbox + worker webhooków; inbound callbacks; CredentialVault (libsodium) na sekrety integracji.",
            "Payment lock w Driver App: dostawa dopiero po rozliczeniu gotówka/karta zgodnie z silnikiem kursów.",
        ],
        "Źródło: _docs/01_KONSTYTUCJA.md, _docs/02_ARCHITEKTURA.md, api/courses/engine.php",
    )

    add_bullets(
        "Perły produktowe — Online Studio",
        [
            "Hollywood Director's Suite: MagicEnhance, MagicBake, MagicRelight, MagicColorGrade, MagicDust, MagicCompanions.",
            "Style Conductor — jednym przebiegiem zmiana tożsamości wizualnej całej kategorii (LUT, światło, companions, typografia).",
            "Harmony Score — transparentny model (kompletność / dopracowanie / spójność) + actionable outliers, nie ozdoba.",
            "Wspólny core/js/scene_renderer.js — WYSIWYG Directora = render storefrontu 1:1 dla warstw pizzy (M2.1).",
            "12 presetów stylów kinowych (m022), LUT Library, kamery (m.in. auto-perspective, presety).",
        ],
        "Źródło: _docs/15_KIERUNEK_ONLINE.md §2.4, _docs/canvasy/SliceHub Three Directions.md",
    )

    add_bullets(
        "Perły operacyjne — co już stoi w kodzie",
        [
            "Warehouse: pełny obieg dokumentów (PZ, RW, KOR, MM, inwentaryzacja) + wiele UI modułu magazynu.",
            "Courses + Dispatcher + Driver PWA: K-system kursów, L-kolejka przystanków, emergency recall (heading=-999).",
            "Track storefrontu: live ETA, hero image, pozycje, pinezka restauracji i kierowcy (_docs/00 — Faza E DONE).",
            "Offline-first POS: świadomy FREEZE do Anti-Corruption Layer — dojrzała decyzja architektoniczna, nie chaos.",
        ],
        "Źródło: _docs/02_ARCHITEKTURA.md, _docs/00_PAMIEC_SYSTEMU.md (FREEZE NOTICE)",
    )

    add_bullets(
        "Innowacja roadmapy (nie slid na „jutro”)",
        [
            "MVP Fazy 2: Droga B — Drzwi + Counter + Living Table (6–8 tyg. ramy z dokumentu decyzyjnego).",
            "Faza 3: Restaurant Viewfinder — jedna scena, swipe w 4 kierunkach (menu / kuchnia / sala / koszyk); unikalny kompromis A/C.",
            "Faza 5+: pełny „Film” pięciu scen — świadomie odłożony (koszt, performance, ryzyko).",
            "G3 AI Jobs (sh_ai_jobs + worker) — odłożone do osobnej fazy AI (jasna granica scope).",
        ],
        "Źródło: _docs/15_KIERUNEK_ONLINE.md, SliceHub Three Directions.md",
    )

    if "online_wide" in png:
        add_image_slide(
            "Storefront PWA — mock wide (1280×720) z repo",
            png["online_wide"],
            "Plik źródłowy: modules/online/screenshots/wide.svg — materiał do wniosku / presskit.",
        )
    if "online_narrow" in png:
        add_image_slide(
            "Storefront PWA — mock mobile (390×844) z repo",
            png["online_narrow"],
            "Plik źródłowy: modules/online/screenshots/narrow.svg",
        )
    if "pos_wide" in png:
        add_image_slide(
            "POS PWA — mock wide z repo",
            png["pos_wide"],
            "Plik źródłowy: modules/pos/screenshots/wide.svg — dark battlefield theme.",
        )
    if "pos_narrow" in png:
        add_image_slide(
            "POS PWA — mock mobile z repo",
            png["pos_narrow"],
            "Plik źródłowy: modules/pos/screenshots/narrow.svg",
        )

    add_bullets(
        "Jak odpalić demo (hosting / lokalnie)",
        [
            "Migracje + setup: database/README.md, INSTRUKCJA_CZYSTY_START.md.",
            "Seed pełnego tenantu demo: scripts/seed_demo_all.php (jednorazowo — uwaga na duplikaty zamówień).",
            "Konta PIN po seedzie: _docs/DEPLOYMENT_HOSTING.md (rotacja PIN przed produkcją).",
        ],
        "Źródło: _docs/DEPLOYMENT_HOSTING.md, database/INSTRUKCJA_CZYSTY_START.md",
    )

    add_bullets(
        "The Ask — Spark 3.0",
        [
            "Środki + mentoring na domknięcie Fazy F (Counter + Living Table) oraz 3 piloty produkcyjne.",
            "KPI propozycja: konwersja doorway→koszyk vs. layout statyczny; Harmony avg per kategoria; brak P0 cross-tenant.",
            "Materiały źródłowe decku: _docs/pitchdeck_SPARK_3_ACCELERATOR.md + ten PPTX/PDF (generowane skryptem).",
        ],
        "Uzupełnij zespół, ARPU, CAC — brak w repo, celowo nie wymyślane.",
    )

    add_title_slide("Dziękujemy", "SliceHub Enterprise · kontakt: (uzupełnij) · demo wg DEPLOYMENT_HOSTING.md")

    out = OUT_DIR / "SliceHub_Spark3_Pitch.pptx"
    prs.save(str(out))
    print(f"OK wrote {out}")
    return out


def build_pdf(png: dict[str, Path]) -> Path:
    from fpdf import FPDF

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    font = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    font_b = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    key_to_src = {k: v for k, v in SVG_SOURCES}

    class Deck(FPDF):
        def footer(self):
            self.set_y(-12)
            self.set_font("DejaVu", "", 8)
            self.set_text_color(120, 120, 120)
            self.cell(0, 8, f"SliceHub Enterprise · strona {self.page_no()}", align="C")

    pdf = Deck(format="A4")
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(14, 16, 14)
    pdf.add_font("DejaVu", "", font)
    pdf.add_font("DejaVu", "B", font_b)
    col_w = pdf.w - pdf.l_margin - pdf.r_margin

    def pg_title(t: str) -> None:
        pdf.add_page()
        pdf.set_x(pdf.l_margin)
        pdf.set_font("DejaVu", "B", 18)
        pdf.set_text_color(40, 32, 24)
        pdf.multi_cell(col_w, 9, t)
        pdf.ln(3)
        pdf.set_x(pdf.l_margin)
        pdf.set_font("DejaVu", "", 11)
        pdf.set_text_color(35, 35, 40)

    pg_title("SliceHub Enterprise — Spark 3.0 Accelerator")
    pdf.set_x(pdf.l_margin)
    pdf.multi_cell(
        col_w,
        6,
        "PDF: scripts/build_spark_pitchdeck.py. Źródła treści: _docs/01_KONSTYTUCJA.md, "
        "_docs/00_PAMIEC_SYSTEMU.md, _docs/15_KIERUNEK_ONLINE.md, "
        "_docs/canvasy/SliceHub Three Directions.md, _docs/02_ARCHITEKTURA.md. "
        "W repozytorium nie ma gotowych PNG/MP4 — użyto mocków SVG PWA (modules/*/screenshots/*.svg) "
        "po rasteryzacji do PNG.",
    )

    pg_title("Problem")
    for line in [
        "Lokal traci markę w agregatorze i spójność między magazynem a tym, co klient widzi online.",
        "Typowy sklep internetowy SKU nie buduje „teatru restauracji”.",
    ]:
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(col_w, 6, "- " + line)

    pg_title("Rozwiązanie — jednym zdaniem")
    pdf.set_x(pdf.l_margin)
    pdf.multi_cell(
        col_w,
        6,
        "Dwie twarze jednego OS: manager w kinowym Studio (Director, Style Conductor, Harmony), "
        "klient w teatrze fotograficznym storefrontu (Drzwi, Living Table, Track) — _docs/00_PAMIEC_SYSTEMU.md.",
    )

    pg_title("Perły — architektura i bezpieczeństwo")
    for line in [
        "CartEngine na serwerze — frontend nie wysyła cen ani totali (Prawo IV).",
        "Multi-tenant: tenant_id; silosy sh_/sys_/wh_; mosty tylko po SKU / ascii_key (Konstytucja, rozdz. 9).",
        "CredentialVault (libsodium), event outbox, inbound adapters.",
        "Driver: payment lock przed statusem „dostarczono”; emergency recall w logistyce.",
    ]:
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(col_w, 6, "- " + line)

    pg_title("Perły — Online Studio / Storefront")
    for line in [
        "Magic: Bake, Relight, ColorGrade, Dust, Companions, Enhance („THE button”).",
        "Style Conductor — tożsamość wizualna całej kategorii jednym przebiegiem.",
        "Harmony Score — metryki actionable (kompletność / dopracowanie / spójność).",
        "core/js/scene_renderer.js — SSOT warstwy: Director i Storefront ten sam rendering pizzy (M2.1).",
        "12 stylów kinowych (m022), LUT Library, presety kamery + auto-perspective.",
        "Roadmapa produktowa: B (Counter+Drzwi) → C (Restaurant Viewfinder) → A (Film) — świadome odłożenie A.",
    ]:
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(col_w, 6, "- " + line)

    pg_title("Perły — operacje")
    for line in [
        "Warehouse: PZ, RW, KOR, MM, inwentaryzacja; receptury BOM → food cost.",
        "Courses + Dispatcher + Driver PWA; K-system, L-przystanki.",
        "Track storefrontu: ETA, hero, pozycje, mapa kierowcy (faza E zamknięta w dokumentacji).",
        "Zero-Reload Runtime: Vanilla JS + PHP na produkcji (_docs/01_KONSTYTUCJA.md).",
    ]:
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(col_w, 6, "- " + line)

    for label, key in [
        ("Storefront — mock wide (z repo)", "online_wide"),
        ("Storefront — mock mobile (z repo)", "online_narrow"),
        ("POS — mock wide (z repo)", "pos_wide"),
        ("POS — mock mobile (z repo)", "pos_narrow"),
    ]:
        pth = png.get(key)
        if not pth or not pth.exists():
            continue
        pdf.add_page()
        pdf.set_x(pdf.l_margin)
        pdf.set_font("DejaVu", "B", 14)
        pdf.set_text_color(40, 32, 24)
        pdf.cell(0, 8, label)
        pdf.ln(10)
        pdf.image(str(pth), x=10, w=190)
        pdf.ln(6)
        pdf.set_font("DejaVu", "", 9)
        pdf.set_text_color(80, 75, 70)
        rel = key_to_src.get(key, Path("."))
        try:
            rels = str(rel.relative_to(ROOT))
        except ValueError:
            rels = str(rel)
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(col_w, 5, f"Źródło: {rels} → PNG (CairoSVG, scale=2).")

    pg_title("Demo w repo")
    pdf.set_x(pdf.l_margin)
    pdf.multi_cell(
        col_w,
        6,
        "Instrukcja: _docs/DEPLOYMENT_HOSTING.md, database/INSTRUKCJA_CZYSTY_START.md. "
        "Seed: scripts/seed_demo_all.php (jednorazowo). Brak nagrania wideo w git — plan: _docs/17_OFFLINE_POS_BACKLOG.md (P8 demo).",
    )

    pg_title("The Ask — Spark 3.0")
    pdf.set_x(pdf.l_margin)
    pdf.multi_cell(
        col_w,
        6,
        "Kapitał / grant na Fazę F (Counter + Living Table) oraz piloty; mentoring pod Prawo VII. "
        "Uzupełnij zespół, ARPU, CAC — brak w repo.",
    )

    out = OUT_DIR / "SliceHub_Spark3_Pitch.pdf"
    pdf.output(str(out))
    print(f"OK wrote {out}")
    return out


def main() -> int:
    png = rasterize_svgs()
    if not png:
        print("ERROR: no PNG generated", file=sys.stderr)
        return 1
    build_pptx(png)
    build_pdf(png)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
